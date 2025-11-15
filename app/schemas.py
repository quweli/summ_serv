from io import BytesIO
from docx import Document
from pptx import Presentation
from openai import OpenAI
from pydantic import BaseModel


class Extractor:
    def __init__(self):
        self.ext = {"txt", "docx", "pptx"}

    def extens(self, filename: str, file_bytes: bytes) -> str:
        ext = filename.rsplit(".", 1)[1].lower()
        if ext not in self.ext:
            raise ValueError(f"Расширение не поддерживается: {ext}")
        if ext == "txt":
            return self.extract_txt(file_bytes)
        if ext == "docx":
            return self.extract_docx(file_bytes)
        if ext == "pptx":
            return self.extract_pptx(file_bytes)   
    def extract_txt(self, data: bytes) -> str:
        try:
            return data.decode("utf-8", errors="ignore")
        except Exception:
            return data.decode("latin-1", errors="ignore")
    def extract_docx(self, data: bytes) -> str:
        f = BytesIO(data)
        doc = Document(f)
        paragraphs = [p.text for p in doc.paragraphs if p.text]
        return "\n".join(paragraphs)
    def extract_pptx(self, data: bytes) -> str:
        f = BytesIO(data)
        prs = Presentation(f)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
        return "\n".join(texts)

class OpenAIClient:
    def __init__(self, model: str):
        self.client = OpenAI()
        self.model = model
        self.prompt = (
            "Ты — профессиональный ассистент по сокращению текста. "
            "Сделай краткое содержание текста (3–6 предложений). "
            "Пиши емко, без лишних слов."
        )
    def summarize(self, text: str, word_count: int | None = None) -> str:
        prompt = self.prompt
        if word_count:
            add = (
                f"Сократи текст до примерно {word_count} слов. "
            )
            prompt += add

        messages = [
            {"role": "system", "content": prompt},
            {"role": "user", "content": f"Суммаризуй следующий текст:\n\n{text}"}
        ]

        response = self.client.chat.completions.create(
            model=self.model,
            messages=messages,
            temperature=0.3,
            max_tokens=700,
        )
        return response.choices[0].message.content.strip()
    
class SummaryResponse(BaseModel):
    summary: str